"""
Módulo unificado de ingestão — Fito+ Amazônia / AIMM.

Suporta ingestão de:
  Arquivos  : xlsx, xls, ods, pptx, docx, pdf, json, yaml, markdown, txt
              + formatos GIS: geojson, shp, gpkg, kml, gml, csv, tsv
  URLs      : download de páginas web ou arquivos remotos
  DOIs      : resolução via API CrossRef (api.crossref.org)
  APIs      : conectores genéricos por endpoint HTTP

Uso rápido:
    from fito_aimm.ingestor import ingerir

    resultado = ingerir("relatorio.pdf")
    resultado = ingerir("https://exemplo.com/dados.xlsx")
    resultado = ingerir("10.1016/j.forpol.2021.102447", tipo="doi")
    resultado = ingerir("https://api.ibge.gov.br/v3/...", tipo="api")
"""
from __future__ import annotations

import hashlib
import json
import re
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import requests

# ---------------------------------------------------------------------------
# Importações opcionais — os parsers degradam graciosamente se ausentes
# ---------------------------------------------------------------------------

try:
    import openpyxl
    _OPENPYXL_OK = True
except ImportError:
    _OPENPYXL_OK = False

try:
    from pypdf import PdfReader
    _PYPDF_OK = True
except ImportError:
    _PYPDF_OK = False

try:
    from docx import Document as DocxDocument
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    from pptx import Presentation
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import pandas as pd
    _PANDAS_OK = True
except ImportError:
    _PANDAS_OK = False

try:
    import yaml
    _YAML_OK = True
except ImportError:
    _YAML_OK = False

try:
    from bs4 import BeautifulSoup
    _BS4_OK = True
except ImportError:
    _BS4_OK = False

try:
    import geopandas as gpd
    _GEOPANDAS_OK = True
except ImportError:
    _GEOPANDAS_OK = False


# ---------------------------------------------------------------------------
# Mapeamento de extensões → categoria de formato
# ---------------------------------------------------------------------------

EXTENSOES_PLANILHA = {".xlsx", ".xls", ".ods"}
EXTENSOES_APRESENTACAO = {".pptx", ".ppt"}
EXTENSOES_DOCUMENTO = {".docx", ".doc"}
EXTENSOES_PDF = {".pdf"}
EXTENSOES_JSON = {".json"}
EXTENSOES_YAML = {".yaml", ".yml"}
EXTENSOES_MARKDOWN = {".md", ".markdown"}
EXTENSOES_TEXTO = {".txt"}
EXTENSOES_TABELA = {".csv", ".tsv"}
EXTENSOES_GIS_NATIVO = {".geojson"}
EXTENSOES_GIS_BINARIO = {".shp", ".gpkg", ".kml", ".gml"}
EXTENSOES_GIS = EXTENSOES_GIS_NATIVO | EXTENSOES_GIS_BINARIO

TODAS_EXTENSOES = (
    EXTENSOES_PLANILHA
    | EXTENSOES_APRESENTACAO
    | EXTENSOES_DOCUMENTO
    | EXTENSOES_PDF
    | EXTENSOES_JSON
    | EXTENSOES_YAML
    | EXTENSOES_MARKDOWN
    | EXTENSOES_TEXTO
    | EXTENSOES_TABELA
    | EXTENSOES_GIS
)

# Padrão para reconhecer DOIs
_REGEX_DOI = re.compile(
    r"^(?:(?:https?://(?:dx\.)?doi\.org/)|(?:doi:))?10\.\d{4,}[^\s]*$",
    re.IGNORECASE,
)

_TIMEOUT_REQUESTS = 30  # segundos


# ---------------------------------------------------------------------------
# Estrutura de resultado
# ---------------------------------------------------------------------------

@dataclass
class ResultadoIngestao:
    """Resultado padronizado de uma operação de ingestão."""

    fonte: str
    tipo_fonte: str          # "arquivo" | "url" | "doi" | "api"
    formato: str             # extensão ou "url" / "doi" / "api"
    titulo: str
    conteudo: dict[str, Any] = field(default_factory=dict)
    metadados: dict[str, Any] = field(default_factory=dict)
    status: str = "sucesso"  # "sucesso" | "parcial" | "erro"
    erro: str | None = None

    # ------------------------------------------------------------------
    # Helpers de acesso rápido
    # ------------------------------------------------------------------

    @property
    def texto(self) -> str:
        """Texto principal extraído (se disponível)."""
        return str(self.conteudo.get("texto", ""))

    @property
    def linhas(self) -> list[dict[str, Any]]:
        """Linhas/registros tabulares (planilhas, CSV, etc.)."""
        return self.conteudo.get("linhas", [])

    @property
    def resumo(self) -> str:
        n_linhas = len(self.linhas)
        texto = self.texto
        trecho = (texto[:200] + "…") if len(texto) > 200 else texto
        if n_linhas:
            return f"{n_linhas} linha(s) | {trecho}"
        return trecho or "(sem conteúdo)"


# ---------------------------------------------------------------------------
# Detecção automática do tipo de fonte
# ---------------------------------------------------------------------------

def _detectar_tipo(fonte: str) -> tuple[str, str]:
    """Retorna (tipo_fonte, formato).

    tipo_fonte: "arquivo" | "url" | "doi" | "api"
    formato:    extensão sem ponto ou tipo genérico ("url", "doi", "api")
    """
    fonte = fonte.strip()

    # DOI explícito
    if _REGEX_DOI.match(fonte):
        return "doi", "doi"

    # URL
    parsed = urlparse(fonte)
    if parsed.scheme in ("http", "https"):
        # Verifica se a URL parece um endpoint de API (sem extensão conhecida)
        path = parsed.path.rstrip("/")
        ext = Path(path).suffix.lower() if path else ""
        if ext in TODAS_EXTENSOES:
            return "url", ext.lstrip(".")
        return "url", "url"

    # Arquivo local
    caminho = Path(fonte)
    ext = caminho.suffix.lower()
    if ext in TODAS_EXTENSOES:
        return "arquivo", ext.lstrip(".")
    return "arquivo", ext.lstrip(".") or "desconhecido"


# ---------------------------------------------------------------------------
# Parsers de arquivo
# ---------------------------------------------------------------------------

def _parsear_planilha(caminho: Path) -> dict[str, Any]:
    if not _OPENPYXL_OK and not _PANDAS_OK:
        return {"erro": "openpyxl ou pandas não disponível"}

    resultado: dict[str, Any] = {"abas": []}
    try:
        if _PANDAS_OK:
            sheets = pd.read_excel(caminho, sheet_name=None)
            for nome_aba, df in sheets.items():
                resultado["abas"].append({
                    "nome": nome_aba,
                    "colunas": list(df.columns.astype(str)),
                    "linhas": df.fillna("").astype(str).to_dict(orient="records"),
                    "total_linhas": len(df),
                })
            # Flatten para compatibilidade: usa a primeira aba como "linhas" raiz
            if resultado["abas"]:
                resultado["linhas"] = resultado["abas"][0]["linhas"]
                resultado["colunas"] = resultado["abas"][0]["colunas"]
        elif _OPENPYXL_OK:
            wb = openpyxl.load_workbook(caminho, read_only=True, data_only=True)
            for ws in wb.worksheets:
                linhas = [
                    {str(ws.cell(1, c).value or f"col_{c}"): str(row[c - 1].value or "")
                     for c in range(1, ws.max_column + 1)}
                    for row in ws.iter_rows(min_row=2, values_only=False)
                ]
                resultado["abas"].append({
                    "nome": ws.title,
                    "linhas": linhas,
                    "total_linhas": len(linhas),
                })
            if resultado["abas"]:
                resultado["linhas"] = resultado["abas"][0]["linhas"]
    except Exception as exc:
        resultado["erro"] = str(exc)
    return resultado


def _parsear_pdf(caminho: Path) -> dict[str, Any]:
    if not _PYPDF_OK:
        return {"erro": "pypdf não disponível. Execute: pip install pypdf"}
    try:
        reader = PdfReader(str(caminho))
        paginas = []
        for i, page in enumerate(reader.pages):
            texto = page.extract_text() or ""
            paginas.append({"pagina": i + 1, "texto": texto})
        texto_completo = "\n\n".join(p["texto"] for p in paginas)
        return {
            "texto": texto_completo,
            "total_paginas": len(paginas),
            "paginas": paginas,
            "metadados_pdf": dict(reader.metadata or {}),
        }
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_docx(caminho: Path) -> dict[str, Any]:
    if not _DOCX_OK:
        return {"erro": "python-docx não disponível. Execute: pip install python-docx"}
    try:
        doc = DocxDocument(str(caminho))
        paragrafos = [p.text for p in doc.paragraphs if p.text.strip()]
        texto = "\n\n".join(paragrafos)
        tabelas = []
        for i, tabela in enumerate(doc.tables):
            linhas_tab = []
            headers: list[str] = []
            for j, row in enumerate(tabela.rows):
                cells = [c.text.strip() for c in row.cells]
                if j == 0:
                    headers = cells
                else:
                    linhas_tab.append(dict(zip(headers, cells)))
            tabelas.append({"tabela": i + 1, "cabecalhos": headers, "linhas": linhas_tab})
        return {
            "texto": texto,
            "paragrafos": paragrafos,
            "tabelas": tabelas,
            "total_paragrafos": len(paragrafos),
        }
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_pptx(caminho: Path) -> dict[str, Any]:
    if not _PPTX_OK:
        return {"erro": "python-pptx não disponível. Execute: pip install python-pptx"}
    try:
        prs = Presentation(str(caminho))
        slides = []
        texto_completo_partes: list[str] = []
        for i, slide in enumerate(prs.slides):
            textos_slide: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        linha = para.text.strip()
                        if linha:
                            textos_slide.append(linha)
            slides.append({"slide": i + 1, "textos": textos_slide})
            texto_completo_partes.extend(textos_slide)
        return {
            "texto": "\n".join(texto_completo_partes),
            "total_slides": len(slides),
            "slides": slides,
        }
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_json(caminho: Path) -> dict[str, Any]:
    try:
        with caminho.open("r", encoding="utf-8") as f:
            dados = json.load(f)
        if isinstance(dados, list):
            return {"linhas": dados, "total_registros": len(dados)}
        if isinstance(dados, dict):
            return {"dados": dados, "chaves": list(dados.keys())}
        return {"dados": dados}
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_yaml(caminho: Path) -> dict[str, Any]:
    if not _YAML_OK:
        return {"erro": "PyYAML não disponível. Execute: pip install PyYAML"}
    try:
        with caminho.open("r", encoding="utf-8") as f:
            dados = yaml.safe_load(f)
        if isinstance(dados, list):
            return {"linhas": dados, "total_registros": len(dados)}
        if isinstance(dados, dict):
            return {"dados": dados, "chaves": list(dados.keys())}
        return {"dados": dados}
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_markdown(caminho: Path) -> dict[str, Any]:
    try:
        texto = caminho.read_text(encoding="utf-8")
        # Extrai títulos (linhas que começam com #)
        titulos = [
            linha.lstrip("#").strip()
            for linha in texto.splitlines()
            if linha.startswith("#")
        ]
        return {
            "texto": texto,
            "titulos": titulos,
            "total_caracteres": len(texto),
            "total_linhas": len(texto.splitlines()),
        }
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_texto(caminho: Path) -> dict[str, Any]:
    try:
        texto = caminho.read_text(encoding="utf-8", errors="replace")
        return {
            "texto": texto,
            "total_caracteres": len(texto),
            "total_linhas": len(texto.splitlines()),
        }
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_tabela(caminho: Path) -> dict[str, Any]:
    if not _PANDAS_OK:
        return {"erro": "pandas não disponível"}
    try:
        sep = "\t" if caminho.suffix.lower() == ".tsv" else ","
        df = pd.read_csv(caminho, sep=sep, dtype=str)
        return {
            "colunas": list(df.columns),
            "linhas": df.fillna("").to_dict(orient="records"),
            "total_linhas": len(df),
        }
    except Exception as exc:
        return {"erro": str(exc)}


def _parsear_gis(caminho: Path) -> dict[str, Any]:
    ext = caminho.suffix.lower()
    if ext == ".geojson":
        try:
            with caminho.open("r", encoding="utf-8") as f:
                geojson = json.load(f)
            features = (
                geojson.get("features", [])
                if geojson.get("type") == "FeatureCollection"
                else [geojson]
            )
            return {
                "tipo_gis": "geojson",
                "total_features": len(features),
                "features": features,
                "propriedades_exemplo": list(features[0].get("properties", {}).keys()) if features else [],
            }
        except Exception as exc:
            return {"erro": str(exc)}

    # Shapefile / GeoPackage / KML / GML — requer geopandas
    if not _GEOPANDAS_OK:
        return {"erro": "geopandas não disponível para este formato GIS. Execute: pip install geopandas"}
    try:
        gdf = gpd.read_file(str(caminho))
        df_sem_geom = gdf.drop(columns=["geometry"], errors="ignore")
        return {
            "tipo_gis": ext.lstrip("."),
            "total_features": len(gdf),
            "crs": str(gdf.crs) if gdf.crs else None,
            "colunas": list(df_sem_geom.columns),
            "linhas": df_sem_geom.fillna("").astype(str).to_dict(orient="records"),
        }
    except Exception as exc:
        return {"erro": str(exc)}


# ---------------------------------------------------------------------------
# Parsers de fontes remotas
# ---------------------------------------------------------------------------

def _parsear_url(url: str, caminho_download: Path | None = None) -> dict[str, Any]:
    """Baixa URL e extrai conteúdo.

    Se a URL terminar com uma extensão de arquivo conhecida, baixa e parseia.
    Caso contrário, extrai texto da página HTML.
    """
    try:
        resp = requests.get(url, timeout=_TIMEOUT_REQUESTS, headers={"User-Agent": "FitoAIMM/1.0"})
        resp.raise_for_status()
    except requests.RequestException as exc:
        return {"erro": str(exc)}

    content_type = resp.headers.get("Content-Type", "")
    # Tenta identificar se é arquivo binário via URL ou Content-Type
    parsed_url = urlparse(url)
    ext = Path(parsed_url.path).suffix.lower()

    if ext in TODAS_EXTENSOES and ext not in {".md", ".json", ".yaml", ".yml", ".txt", ".csv", ".tsv"}:
        # Arquivo binário — salva temporariamente e parseia
        if caminho_download is None:
            nome = Path(parsed_url.path).name or "download"
            caminho_download = Path(f"/tmp/fito_aimm_download_{nome}")
        caminho_download.write_bytes(resp.content)
        resultado = _parsear_arquivo(caminho_download)
        resultado["url_origem"] = url
        resultado["content_type"] = content_type
        return resultado

    # Página HTML ou texto
    if "html" in content_type or ext in ("", ".html", ".htm"):
        if _BS4_OK:
            soup = BeautifulSoup(resp.content, "lxml")
            titulo = soup.find("title")
            titulo_texto = titulo.get_text(strip=True) if titulo else ""
            # Remove scripts e styles
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            texto = soup.get_text(separator="\n", strip=True)
            return {
                "titulo_html": titulo_texto,
                "texto": texto[:50000],  # limita a 50k chars
                "content_type": content_type,
                "status_http": resp.status_code,
            }
        return {"texto": resp.text[:50000], "content_type": content_type}

    # JSON
    if "json" in content_type or ext == ".json":
        try:
            dados = resp.json()
            if isinstance(dados, list):
                return {"linhas": dados, "total_registros": len(dados), "content_type": content_type}
            return {"dados": dados, "chaves": list(dados.keys()) if isinstance(dados, dict) else [], "content_type": content_type}
        except Exception:
            pass

    # Texto simples
    return {"texto": resp.text[:50000], "content_type": content_type, "status_http": resp.status_code}


def _parsear_doi(doi: str) -> dict[str, Any]:
    """Resolve metadados de um DOI via API CrossRef."""
    # Normaliza o DOI (remove prefixo se houver)
    doi_limpo = re.sub(r"^(?:https?://(?:dx\.)?doi\.org/|doi:)", "", doi, flags=re.IGNORECASE).strip()

    url_crossref = f"https://api.crossref.org/works/{doi_limpo}"
    try:
        resp = requests.get(
            url_crossref,
            timeout=_TIMEOUT_REQUESTS,
            headers={"User-Agent": "FitoAIMM/1.0 (mailto:contato@fitomaisamazonia.org)"},
        )
        resp.raise_for_status()
    except requests.RequestException as exc:
        return {"erro": str(exc), "doi": doi_limpo}

    try:
        dados = resp.json()
    except Exception as exc:
        return {"erro": f"Resposta inválida do CrossRef: {exc}", "doi": doi_limpo}

    work = dados.get("message", {})
    autores_raw = work.get("author", [])
    autores = [
        f"{a.get('given', '')} {a.get('family', '')}".strip()
        for a in autores_raw
    ]
    data_pub = work.get("published", {}).get("date-parts", [[None]])[0]
    ano = data_pub[0] if data_pub else None

    titulo_lista = work.get("title", [])
    titulo = titulo_lista[0] if titulo_lista else ""

    abstract_raw = work.get("abstract", "")
    # Remove tags XML do resumo (CrossRef retorna HTML/XML)
    abstract = re.sub(r"<[^>]+>", " ", abstract_raw).strip()

    return {
        "doi": doi_limpo,
        "titulo": titulo,
        "autores": autores,
        "ano_publicacao": ano,
        "publicador": work.get("publisher", ""),
        "tipo_publicacao": work.get("type", ""),
        "resumo": abstract,
        "url": f"https://doi.org/{doi_limpo}",
        "issn": work.get("ISSN", []),
        "container_title": work.get("container-title", []),
        "subject": work.get("subject", []),
    }


def _parsear_api(
    url: str,
    metodo: str = "GET",
    parametros: dict[str, Any] | None = None,
    cabecalhos: dict[str, str] | None = None,
    corpo: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Chama um endpoint de API REST e retorna o resultado."""
    hdrs = {"User-Agent": "FitoAIMM/1.0", "Accept": "application/json"}
    if cabecalhos:
        hdrs.update(cabecalhos)

    try:
        resp = requests.request(
            metodo.upper(),
            url,
            params=parametros,
            headers=hdrs,
            json=corpo,
            timeout=_TIMEOUT_REQUESTS,
        )
        resp.raise_for_status()
    except requests.RequestException as exc:
        return {"erro": str(exc), "url": url}

    content_type = resp.headers.get("Content-Type", "")
    try:
        dados = resp.json()
    except Exception:
        dados = resp.text[:50000]

    return {
        "url": url,
        "metodo": metodo.upper(),
        "status_http": resp.status_code,
        "content_type": content_type,
        "dados": dados,
    }


# ---------------------------------------------------------------------------
# Dispatcher por extensão de arquivo
# ---------------------------------------------------------------------------

def _parsear_arquivo(caminho: Path) -> dict[str, Any]:
    ext = caminho.suffix.lower()
    if ext in EXTENSOES_PLANILHA:
        return _parsear_planilha(caminho)
    if ext in EXTENSOES_PDF:
        return _parsear_pdf(caminho)
    if ext in EXTENSOES_DOCUMENTO:
        return _parsear_docx(caminho)
    if ext in EXTENSOES_APRESENTACAO:
        return _parsear_pptx(caminho)
    if ext in EXTENSOES_JSON:
        return _parsear_json(caminho)
    if ext in EXTENSOES_YAML:
        return _parsear_yaml(caminho)
    if ext in EXTENSOES_MARKDOWN:
        return _parsear_markdown(caminho)
    if ext in EXTENSOES_TEXTO:
        return _parsear_texto(caminho)
    if ext in EXTENSOES_TABELA:
        return _parsear_tabela(caminho)
    if ext in EXTENSOES_GIS:
        return _parsear_gis(caminho)
    return {"erro": f"Formato '{ext}' não suportado pelo ingestor."}


def _titulo_de_arquivo(caminho: Path, conteudo: dict[str, Any]) -> str:
    """Tenta extrair um título significativo do conteúdo."""
    if "titulo" in conteudo and conteudo["titulo"]:
        return str(conteudo["titulo"])
    if "titulo_html" in conteudo and conteudo["titulo_html"]:
        return str(conteudo["titulo_html"])
    metadados_pdf = conteudo.get("metadados_pdf", {})
    if "/Title" in metadados_pdf:
        return str(metadados_pdf["/Title"])
    # Primeira linha de texto como fallback
    texto = conteudo.get("texto", "")
    if texto:
        primeira_linha = texto.splitlines()[0].strip()
        if primeira_linha:
            return primeira_linha[:120]
    return caminho.name


# ---------------------------------------------------------------------------
# API pública principal
# ---------------------------------------------------------------------------

def ingerir(
    fonte: str,
    tipo: str | None = None,
    *,
    # Opções para API
    metodo_api: str = "GET",
    parametros_api: dict[str, Any] | None = None,
    cabecalhos_api: dict[str, str] | None = None,
    corpo_api: dict[str, Any] | None = None,
    # Opções gerais
    salvar_em: Path | None = None,
) -> ResultadoIngestao:
    """Ingere uma fonte de dados e retorna um ResultadoIngestao padronizado.

    Args:
        fonte: Caminho de arquivo, URL, DOI ou endpoint de API.
        tipo: Tipo explícito ("arquivo", "url", "doi", "api"). Se None, é
              detectado automaticamente a partir da fonte.
        metodo_api: Método HTTP para chamadas de API (padrão: "GET").
        parametros_api: Parâmetros de query string para APIs.
        cabecalhos_api: Cabeçalhos HTTP adicionais para APIs.
        corpo_api: Corpo da requisição (JSON) para APIs POST/PUT.
        salvar_em: Diretório onde salvar o resultado em JSON (opcional).

    Returns:
        ResultadoIngestao com status, conteúdo e metadados.
    """
    tipo_detectado, formato_detectado = _detectar_tipo(fonte)
    tipo_final = tipo or tipo_detectado

    titulo = ""
    conteudo: dict[str, Any] = {}
    metadados: dict[str, Any] = {
        "fonte": fonte,
        "tipo_fonte": tipo_final,
        "formato": formato_detectado,
        "timestamp_ingestao": _timestamp_utc(),
    }
    status = "sucesso"
    erro: str | None = None

    try:
        if tipo_final == "arquivo":
            caminho = Path(fonte)
            if not caminho.exists():
                raise FileNotFoundError(f"Arquivo não encontrado: {fonte}")
            conteudo = _parsear_arquivo(caminho)
            titulo = _titulo_de_arquivo(caminho, conteudo)
            metadados["tamanho_bytes"] = caminho.stat().st_size
            metadados["hash_sha256"] = _hash_arquivo(caminho)

        elif tipo_final == "url":
            conteudo = _parsear_url(fonte)
            titulo = str(conteudo.get("titulo_html", "") or urlparse(fonte).path.rstrip("/").split("/")[-1] or fonte)

        elif tipo_final == "doi":
            conteudo = _parsear_doi(fonte)
            titulo = str(conteudo.get("titulo", fonte))

        elif tipo_final == "api":
            conteudo = _parsear_api(
                fonte,
                metodo=metodo_api,
                parametros=parametros_api,
                cabecalhos=cabecalhos_api,
                corpo=corpo_api,
            )
            titulo = f"API: {fonte}"

        else:
            raise ValueError(f"Tipo de fonte desconhecido: '{tipo_final}'")

        if "erro" in conteudo:
            status = "parcial"
            erro = conteudo.pop("erro")

    except Exception as exc:
        status = "erro"
        erro = str(exc)
        conteudo = {}

    resultado = ResultadoIngestao(
        fonte=fonte,
        tipo_fonte=tipo_final,
        formato=formato_detectado,
        titulo=titulo,
        conteudo=conteudo,
        metadados=metadados,
        status=status,
        erro=erro,
    )

    if salvar_em is not None:
        _salvar_resultado(resultado, Path(salvar_em))

    return resultado


def ingerir_lote(
    fontes: list[str | dict[str, Any]],
    salvar_em: Path | None = None,
    pausa_entre_requisicoes: float = 0.5,
) -> list[ResultadoIngestao]:
    """Ingere múltiplas fontes em sequência.

    Args:
        fontes: Lista de fontes. Cada item pode ser:
            - str: caminho, URL, DOI ou endpoint de API
            - dict: {"fonte": "...", "tipo": "...", ...kwargs}
        salvar_em: Diretório onde salvar resultados individuais em JSON.
        pausa_entre_requisicoes: Segundos de pausa entre chamadas remotas.

    Returns:
        Lista de ResultadoIngestao, um por fonte.
    """
    resultados: list[ResultadoIngestao] = []
    for item in fontes:
        if isinstance(item, str):
            kwargs: dict[str, Any] = {"fonte": item}
        else:
            kwargs = dict(item)

        fonte = kwargs.pop("fonte")
        kwargs["salvar_em"] = salvar_em
        resultado = ingerir(fonte, **kwargs)
        resultados.append(resultado)

        # Pausa entre requisições remotas para ser respeitoso com os servidores
        tipo_final = resultado.tipo_fonte
        if tipo_final in ("url", "doi", "api"):
            time.sleep(pausa_entre_requisicoes)

    return resultados


# ---------------------------------------------------------------------------
# Utilitários internos
# ---------------------------------------------------------------------------

def _timestamp_utc() -> str:
    from datetime import datetime, timezone
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")


def _hash_arquivo(caminho: Path, block_size: int = 65536) -> str:
    h = hashlib.sha256()
    with caminho.open("rb") as f:
        for bloco in iter(lambda: f.read(block_size), b""):
            h.update(bloco)
    return h.hexdigest()


def _salvar_resultado(resultado: ResultadoIngestao, diretorio: Path) -> Path:
    diretorio.mkdir(parents=True, exist_ok=True)
    nome = re.sub(r"[^\w\-.]", "_", resultado.fonte)[:80]
    arquivo = diretorio / f"{nome}.json"
    payload = {
        "fonte": resultado.fonte,
        "tipo_fonte": resultado.tipo_fonte,
        "formato": resultado.formato,
        "titulo": resultado.titulo,
        "status": resultado.status,
        "erro": resultado.erro,
        "metadados": resultado.metadados,
        "conteudo": resultado.conteudo,
    }
    arquivo.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return arquivo


# ---------------------------------------------------------------------------
# Listagem de capacidades
# ---------------------------------------------------------------------------

def capacidades() -> dict[str, Any]:
    """Retorna dicionário descrevendo as capacidades do ingestor."""
    return {
        "formatos_arquivo": {
            "planilhas": sorted(EXTENSOES_PLANILHA),
            "apresentacoes": sorted(EXTENSOES_APRESENTACAO),
            "documentos": sorted(EXTENSOES_DOCUMENTO),
            "pdf": sorted(EXTENSOES_PDF),
            "json": sorted(EXTENSOES_JSON),
            "yaml": sorted(EXTENSOES_YAML),
            "markdown": sorted(EXTENSOES_MARKDOWN),
            "texto": sorted(EXTENSOES_TEXTO),
            "tabelas": sorted(EXTENSOES_TABELA),
            "gis": sorted(EXTENSOES_GIS),
        },
        "fontes_remotas": ["url", "doi", "api"],
        "dependencias": {
            "openpyxl (planilhas xlsx)": _OPENPYXL_OK,
            "pypdf (pdf)": _PYPDF_OK,
            "python-docx (docx)": _DOCX_OK,
            "python-pptx (pptx)": _PPTX_OK,
            "pandas (csv/tabelas/xlsx)": _PANDAS_OK,
            "PyYAML (yaml)": _YAML_OK,
            "beautifulsoup4 (html/url)": _BS4_OK,
            "geopandas (shp/gpkg/kml/gml)": _GEOPANDAS_OK,
        },
    }
